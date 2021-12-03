from pptx import Presentation
from datetime import datetime
from pptx.util import Pt
import os
from pptx.util import Cm

title_for_pptx = '岸田総理大臣就任会見'
words_for_pptx = ["総理",
"自由民主党",
"公明党",
"海軍",
"内閣",
"コロナウイルス"]

def delete_slide(prs, slide):
    #Make dictionary with necessary information
    id_dict = { slide.id: [i, slide.rId] for i,slide in enumerate(prs.slides._sldIdLst) }
    slide_id = slide.slide_id
    prs.part.drop_rel(id_dict[slide_id][1])
    del prs.slides._sldIdLst[id_dict[slide_id][0]]

print("Creating PowerPoint slide...")
prs = Presentation("template.pptx")
runpath = os.getcwd()

bullet_slide_layout = prs.slide_layouts[1]  # page layout
slide = prs.slides.add_slide(bullet_slide_layout)

shapes = slide.shapes
title_shape = shapes.title
body_shape = shapes.placeholders[1]

title_shape.text = title_for_pptx
title_shape.text_frame.paragraphs[0].font.size = Pt(40)  # font size
title_shape.text_frame.paragraphs[0].font.bold = True  # font bold

tf = body_shape.text_frame
tf.margin_top = Cm(10)
print(tf.margin_top)

delete_slide(prs,prs.slides[0])

for i in range (0, min(6,len(words_for_pptx))):
    p = tf.add_paragraph()
    p.text = words_for_pptx[i]
    p.font.size = Pt(20)

prs.save('generated_slide.pptx')
print("Saved PowerPoint slide to",os.path.join(runpath,'generated_slide.pptx'))