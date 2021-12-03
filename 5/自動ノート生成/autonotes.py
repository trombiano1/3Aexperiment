#encoding: UTF-8
import os, sys, subprocess
import cv2
import re
import shutil
import pathlib
import argparse
import ffmpeg #音声処理用
import speech_recognition as sr #音声認識
import MeCab #品詞分解/keyBERT
from keybert import KeyBERT #keyBERT
from pptx import Presentation #パワポ
from pptx.util import Cm, Pt

#============================================
# dependencies: ffmpeg, SpeechRecognition, MeCab, keyBERT, PyTorch (どれもpipで入る)
#
# 使い方は
# python3 project.py -h
# で出る
#
# -i 入力音声/動画の相対path(必須, 動画mp4でしか試してない)
# -s 開始時間と終了時間(秒)を指定できる(デフォルトは全部)
# -t をつけると音声認識の文章を保存する
# -k をつけると抽出されたキーワードをプリント
# -l で言語を指定("ja"(デフォルト)か"en")
#
# ex) $ python3 project.py -i kishida.mp4 -t -s 150 350
#       -> kishida.mp4を150秒~350秒で処理, 文章も保存
#============================================

# ffmpegで音声をカット----------------------------
def split_audio(path, time_section):
    directorypath = os.path.join(pathlib.Path(path).parent.resolve(),'project_data')
    if not os.path.exists(directorypath):
        os.makedirs(directorypath)
    else:
        shutil.rmtree(directorypath)           # Removes all the subdirectories!
        os.makedirs(directorypath)

    audio_input = ffmpeg.input(path)
    info=ffmpeg.probe(path)
    original_duration = int(float(info['format']['duration']))
    if time_section == [-1,-1]:
        time_section = [0, original_duration]
    duration = time_section[1] - time_section[0]
    print("Duration: "+str(duration)+"s")

    print("Saving a frame...")
    vidcap = cv2.VideoCapture(str(path))
    vidcap.set(cv2.CAP_PROP_POS_FRAMES, int(30*duration/2)+time_section[0]*30)
    success,image = vidcap.read()
    cv2.imwrite(str(os.path.join(pathlib.Path(path).parent.resolve(),"project_data","frame.jpg")), image)     # save frame as JPEG file      

    print("Splitting audio....")
    paths = []
    for i in range (0,int((duration-1)/170.0)+1):
        print("Saving audio",i+1,"of",int(duration/170.0)+1,"...")
        if i == int(duration/170.0) + 1:
            audio_cut = audio_input.audio.filter('atrim', duration=(duration%170), start = i * 170 + time_section[0])
        else:
            audio_cut = audio_input.audio.filter('atrim', duration=170, start = i * 170 + time_section[0])
        audio_output = ffmpeg.output(audio_cut, os.path.join(directorypath,('sample'+str(i)+'.wav')),loglevel="quiet")
        paths.append(directorypath+'/sample'+str(i)+'.wav')
        ffmpeg.run(audio_output)
    return paths
    
def audios_to_txt(paths, lang):
    all_text = ""
    if lang == "ja":
        thelang = "ja-JP"
    else:
        thelang = "en-US"
    for i in range (0, len(paths)):
        print("Recognizing audio",i+1,"of",len(paths),"....")
        AUDIO_FILE = paths[i]
        # use the audio file as the audio source
        r = sr.Recognizer()
        with sr.AudioFile(AUDIO_FILE) as source:
            audio = r.record(source)  # read the entire audio file

        result=r.recognize_google(audio, language=thelang)
        # try:
        #     print("Google Speech Recognition thinks you said " + result)
        # except sr.UnknownValueError:
        #     print("Google Speech Recognition could not understand audio")
        # except sr.RequestError as e:
        #     print("Could not request results from Google Speech Recognition service; {0}".format(e))
        all_text += " " + result
    return all_text

# MeCab品詞分解--------------------------------------
def modify_tense(input_sentence):
    print("Recognizing PoS....")
    tokenizer = MeCab.Tagger("-Ochasen -d /usr/local/lib/mecab/dic/mecab-ipadic-neologd")
    tokenizer.parse("")
    node = tokenizer.parseToNode(input_sentence)
    keywords = []
    while node: #品詞ごとの処理
        # if node.feature.split(",")[0] == u"名詞":
        #     keywords.append(node.surface)
        # if node.feature.split(",")[0] == u"形容詞":
        #     #keywords.append(node.feature.split(",")[6]) #原形
        #     keywords.append(node.surface)
        # elif node.feature.split(",")[0] == u"動詞":
        #     #keywords.append(node.feature.split(",")[6]) #原形
        #     keywords.append(node.surface)
        # else:
        #     keywords.append(node.surface)
        keywords.append(node.surface)
        node = node.next
    return ''.join(keywords)

def pos_analysis(input_sentence):
    # 品詞分解 
    tokenizer = MeCab.Tagger("-Ochasen -d /usr/local/lib/mecab/dic/mecab-ipadic-neologd")
    #print(tokenizer)
    tokenizer.parse("")
    node = tokenizer.parseToNode(input_sentence)
    while node:
        #単語を取得
        word = node.surface
        #品詞を取得
        pos = node.feature.split(",")[1]
        print('{0} , {1}'.format(word, pos))
        #次の単語に進める
        node = node.next

def pos_filter(input_words_cluster,lang):
    # 品詞分解 
    tokenizer = MeCab.Tagger("-Ochasen -d /usr/local/lib/mecab/dic/mecab-ipadic-neologd")
    #print(tokenizer)
    tokenizer.parse("")
    node = tokenizer.parseToNode(input_words_cluster)
    return_words = []
    ng_list = ['自立','語幹','接尾','*','助詞','数','接尾']
    while node:
        add_flag = 1
        #単語を取得
        word = node.surface
        #品詞を取得
        pos = node.feature.split(",")[1]
        #print('{0} , {1}'.format(word, pos))
        #次の単語に進める
        if lang == "ja":
            for j in range (0, len(ng_list)):
                if (ng_list[j] in pos) or re.search('[a-zA-Z]', word):
                    add_flag = 0
                    #print(ng_list[j] in pos, any(c.isalpha() for c in word))
        else:
            for j in range (0, len(ng_list)):
                if (ng_list[j] in pos):
                    add_flag = 0
                    #print(ng_list[j] in pos, any(c.isalpha() for c in word))
        if add_flag:
            return_words.append(word)
        node = node.next
    return return_words

# keyBERT で抽出 --------------------------------------
def keyword_extraction(input_sentence, lang):
    print("Extracting keywords...")
    # MeCabで分かち書き
    if lang == "ja":
        words = MeCab.Tagger("-Owakati").parse(input_sentence)
        model = KeyBERT("paraphrase-multilingual-MiniLM-L12-v2"); #for multilingual
    else:
        words = input_sentence
        model = KeyBERT("all-MiniLM-L6-v2") #For english documents
    return model.extract_keywords(words, top_n = 30, keyphrase_ngram_range=(1, 1),use_mmr=True, diversity=0.5)

# Create Presentation ---------------------------------
def delete_slide(prs, slide):
    #Make dictionary with necessary information
    id_dict = { slide.id: [i, slide.rId] for i,slide in enumerate(prs.slides._sldIdLst) }
    slide_id = slide.slide_id
    prs.part.drop_rel(id_dict[slide_id][1])
    del prs.slides._sldIdLst[id_dict[slide_id][0]]

def create_presentation(path, title_for_pptx,words_for_pptx):
    print("Creating PowerPoint slide...")
    directorypath = pathlib.Path(path).parent.resolve()
    prs = Presentation(os.path.join(os.path.dirname(os.path.abspath(__file__)),"template.pptx"))
    runpath = os.getcwd()

    bullet_slide_layout = prs.slide_layouts[1]  # page layout
    slide = prs.slides.add_slide(bullet_slide_layout)

    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]

    title_shape.text = title_for_pptx
    title_shape.text_frame.paragraphs[0].font.size = Pt(40)  # font size
    title_shape.text_frame.paragraphs[0].font.bold = True  # font bold

    tf = body_shape.text_frame
    tf.text = words_for_pptx[0]
    tf.paragraphs[0].font.size = Pt(35)
    print(tf.margin_top)
    
    delete_slide(prs,prs.slides[0])

    for i in range (1, min(5,len(words_for_pptx))):
        p = tf.add_paragraph()
        p.text = words_for_pptx[i]
        p.font.size = Pt(35)
    
    shapes.add_picture(str(os.path.join(directorypath,"project_data",'frame.jpg')), Cm(20), Cm(11), width=Cm(10))

    prs.save(os.path.join(directorypath,'generated_slide.pptx'))
    print("Saved PowerPoint slide to",os.path.join(directorypath,'generated_slide.pptx'))
    return os.path.join(runpath,os.path.join(directorypath,'generated_slide.pptx'))

def open_file(filename):
    if sys.platform == "win32":
        os.startfile(filename)
    else:
        opener = "open" if sys.platform == "darwin" else "xdg-open"
        subprocess.call([opener, filename])

# parse-----------------------
parser = argparse.ArgumentParser()
parser.add_argument('-i', metavar='INPUT_FILE_PATH',action='store', type=pathlib.Path, dest='input_path', help='input file path',required=True)
parser.add_argument('-l', metavar='LANGUAGE',action='store', type=str, dest='lang', help='language', default="ja")
parser.add_argument('-t', action='store_true', dest='save_transcript')
parser.add_argument('-k', action='store_true', dest='print_keywords')
parser.add_argument('-s', metavar=['START_SECOND','END_SECOND'], action='store', type=int, nargs=2, dest='time_section')
args = parser.parse_args()
input_path = args.input_path
title_for_pptx = os.path.basename(input_path).split(".")[0];
#n_soundbites = int(args.n_soundbites)
save_transcript = args.save_transcript
print_keywords = args.print_keywords
lang = args.lang
if args.time_section == None:
    time_section = [-1,-1]
else:
    time_section = args.time_section

if time_section[0] > time_section[1]:
    exit(1)

#main---------------------------------------
paths = split_audio(input_path, time_section)
speech_recognition_text = audios_to_txt(paths, lang)

#parse--------------------------------------
if save_transcript:
    f = open(os.path.join(pathlib.Path(input_path).parent.resolve(),"project_data","TRANSCRIPT.txt"), "x")
    f.write(speech_recognition_text)
    f.close()
    print("Saved transcript to",os.path.join(pathlib.Path(input_path).parent.resolve(),"project_data","TRANSCRIPT.txt"))

if lang == "ja":
    tense_modified_text = modify_tense(speech_recognition_text)
else:
    tense_modified_text = speech_recognition_text


extracted_keywords = keyword_extraction(tense_modified_text,lang)


keywords = []
for i in range (0, len(extracted_keywords)):
    if not(extracted_keywords[i][0] in keywords):
        keywords.append(extracted_keywords[i][0])

sentence = ""
for i in range (0, len(keywords)):
    sentence += keywords[i] + " "

pos_analysis(sentence)


words_for_pptx = pos_filter(sentence,lang)

pptx_list = []
for i in words_for_pptx:
    if not(i in pptx_list):
        pptx_list.append(i)

pptpath = create_presentation(input_path,title_for_pptx, pptx_list)
open_file(pptpath) #throws error, but fine.